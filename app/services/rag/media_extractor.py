"""
Media Extraction Service for RAG.

Extracts text and descriptions from various file types:
- Images: Vision LLM describes content (objects, text, scenes)
- Videos: Extract key frames, describe each with vision LLM
- Excel/CSV: Parse spreadsheets into structured text
- Word/PowerPoint: Extract text content
"""

import os
import base64
import tempfile
from pathlib import Path
from typing import Optional, List, Dict, Any

# Image processing
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

# Video processing (ffmpeg)
try:
    import subprocess
    FFMPEG_AVAILABLE = subprocess.run(
        ["ffmpeg", "-version"], capture_output=True
    ).returncode == 0
except Exception:
    FFMPEG_AVAILABLE = False

# Excel/CSV processing
try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import csv
    CSV_AVAILABLE = True
except ImportError:
    CSV_AVAILABLE = True  # Built-in

# Word documents
try:
    from docx import Document as DocxDocument
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

# PowerPoint
try:
    from pptx import Presentation
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False


# Supported file type groups
IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.tiff', '.tif'}
VIDEO_EXTENSIONS = {'.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv', '.wmv'}
EXCEL_EXTENSIONS = {'.xlsx', '.xls'}
CSV_EXTENSIONS = {'.csv'}
WORD_EXTENSIONS = {'.docx'}
PPT_EXTENSIONS = {'.pptx'}


class MediaExtractorService:
    """
    Extracts text and descriptions from media files for RAG vectorization.
    
    For images and videos, uses a configured vision LLM to generate
    descriptions that are then vectorized as searchable text.
    """

    def __init__(self, vision_provider=None):
        """
        Initialize the media extractor.
        
        Args:
            vision_provider: An LLM provider with vision capability
                            (e.g., OpenAI GPT-4o, Google Gemini, MiMo Omni)
        """
        self._vision_provider = vision_provider

    def set_vision_provider(self, provider):
        """Set or update the vision LLM provider."""
        self._vision_provider = provider

    def can_extract(self, file_path: Path) -> bool:
        """Check if this extractor can handle the given file."""
        suffix = file_path.suffix.lower()
        return suffix in (
            IMAGE_EXTENSIONS | VIDEO_EXTENSIONS | EXCEL_EXTENSIONS |
            CSV_EXTENSIONS | WORD_EXTENSIONS | PPT_EXTENSIONS
        )

    def extract(self, file_path: Path) -> List[Dict[str, Any]]:
        """
        Extract text/description chunks from a media file.
        
        Returns:
            List of dicts with 'content' (str) and 'metadata' (dict)
        """
        suffix = file_path.suffix.lower()

        if suffix in IMAGE_EXTENSIONS:
            return self._extract_image(file_path)
        elif suffix in VIDEO_EXTENSIONS:
            return self._extract_video(file_path)
        elif suffix in EXCEL_EXTENSIONS:
            return self._extract_excel(file_path)
        elif suffix in CSV_EXTENSIONS:
            return self._extract_csv(file_path)
        elif suffix in WORD_EXTENSIONS:
            return self._extract_docx(file_path)
        elif suffix in PPT_EXTENSIONS:
            return self._extract_pptx(file_path)
        return []

    # ============================================================
    # Image Extraction
    # ============================================================
    def _extract_image(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract description from an image using vision LLM."""
        try:
            # Read and encode image
            with open(file_path, "rb") as f:
                image_data = f.read()
            b64_image = base64.b64encode(image_data).decode("utf-8")
            suffix = file_path.suffix.lower()
            mime = {
                '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
                '.png': 'image/png', '.gif': 'image/gif',
                '.webp': 'image/webp', '.bmp': 'image/bmp',
                '.tiff': 'image/tiff', '.tif': 'image/tiff'
            }.get(suffix, 'image/jpeg')

            # Get image dimensions if PIL available
            dimensions = ""
            if PIL_AVAILABLE:
                try:
                    img = Image.open(file_path)
                    dimensions = f" ({img.width}x{img.height}px)"
                except Exception:
                    pass

            description = self._describe_image_with_llm(b64_image, mime)

            if not description:
                provider_name = self._vision_provider.provider_name if self._vision_provider else "none"
                description = (
                    f"[Image file: {file_path.name}{dimensions}] "
                    f"Vision model not available (provider: {provider_name}). "
                    f"To enable image understanding, configure a vision-capable engine "
                    f"(GPT-4o, Gemini 3 Flash/Pro, MiMo Omni, or Claude) with an API key in Settings > Engine Bank."
                )

            return [{
                "content": f"[IMAGE: {file_path.name}{dimensions}]\n{description}",
                "metadata": {
                    "source": str(file_path),
                    "media_type": "image",
                    "file_name": file_path.name,
                    "chunk_index": 0
                }
            }]
        except Exception as e:
            print(f"[MediaExtractor] Image extraction failed for {file_path}: {e}")
            return [{
                "content": f"[IMAGE: {file_path.name}] - Extraction failed: {str(e)}",
                "metadata": {"source": str(file_path), "media_type": "image", "error": str(e)}
            }]

    def _describe_image_with_llm(self, b64_image: str, mime_type: str) -> Optional[str]:
        """Use the vision LLM to describe an image."""
        if not self._vision_provider:
            print("[MediaExtractor] No vision provider configured")
            return None

        try:
            provider = self._vision_provider
            # Ensure client is initialized
            if provider._client is None:
                provider._initialize_client()

            if not hasattr(provider, '_client') or provider._client is None:
                print("[MediaExtractor] Vision provider client failed to initialize")
                return None

            response = provider._client.chat.completions.create(
                model=provider.model_id,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": (
                            "Describe this image in detail for search indexing. "
                            "Include: objects, people, text visible, scene, colors, "
                            "any data or charts shown, and the overall context. "
                            "Be thorough but concise (200-500 words)."
                        )},
                        {"type": "image_url", "image_url": {
                            "url": f"data:{mime_type};base64,{b64_image}"
                        }}
                    ]
                }],
                max_tokens=1000
            )
            result = response.choices[0].message.content
            if result:
                print(f"[MediaExtractor] Vision description generated ({len(result)} chars)")
            return result
        except Exception as e:
            print(f"[MediaExtractor] Vision LLM error: {type(e).__name__}: {e}")
        return None

    # ============================================================
    # Video Extraction
    # ============================================================
    def _extract_video(self, file_path: Path) -> List[Dict[str, Any]]:
        """
        Extract key frames from video and describe each with vision LLM.
        Also extracts audio transcript if possible.
        """
        chunks = []

        if not FFMPEG_AVAILABLE:
            return [{
                "content": f"[VIDEO: {file_path.name}] - ffmpeg not available for frame extraction.",
                "metadata": {"source": str(file_path), "media_type": "video", "error": "ffmpeg_missing"}
            }]

        try:
            # Get video metadata
            duration = self._get_video_duration(file_path)

            # Extract key frames (1 frame per N seconds, max 10 frames)
            frames = self._extract_video_frames(file_path, max_frames=10)

            if not frames:
                return [{
                    "content": f"[VIDEO: {file_path.name}] - Could not extract frames.",
                    "metadata": {"source": str(file_path), "media_type": "video"}
                }]

            # Describe each frame
            for i, (frame_path, timestamp) in enumerate(frames):
                with open(frame_path, "rb") as f:
                    b64 = base64.b64encode(f.read()).decode("utf-8")

                description = self._describe_image_with_llm(b64, "image/jpeg")
                if not description:
                    description = f"[Video frame at {timestamp}s]"

                chunks.append({
                    "content": f"[VIDEO: {file_path.name} | Frame at {timestamp:.1f}s | Duration: {duration:.1f}s]\n{description}",
                    "metadata": {
                        "source": str(file_path),
                        "media_type": "video",
                        "frame_index": i,
                        "timestamp": timestamp,
                        "duration": duration,
                        "file_name": file_path.name
                    }
                })

                # Cleanup temp frame
                try:
                    os.unlink(frame_path)
                except Exception:
                    pass

            # Also try to extract audio transcript
            transcript = self._extract_audio_transcript(file_path)
            if transcript:
                chunks.append({
                    "content": f"[VIDEO: {file_path.name} | Audio Transcript | Duration: {duration:.1f}s]\n{transcript}",
                    "metadata": {
                        "source": str(file_path),
                        "media_type": "video_transcript",
                        "file_name": file_path.name,
                        "duration": duration
                    }
                })

        except Exception as e:
            print(f"[MediaExtractor] Video extraction failed for {file_path}: {e}")
            chunks.append({
                "content": f"[VIDEO: {file_path.name}] - Extraction failed: {str(e)}",
                "metadata": {"source": str(file_path), "media_type": "video", "error": str(e)}
            })

        return chunks if chunks else [{
            "content": f"[VIDEO: {file_path.name}] - No content extracted.",
            "metadata": {"source": str(file_path), "media_type": "video"}
        }]

    def _get_video_duration(self, file_path: Path) -> float:
        """Get video duration in seconds using ffprobe."""
        try:
            result = subprocess.run(
                ["ffprobe", "-v", "error", "-show_entries", "format=duration",
                 "-of", "default=noprint_wrappers=1:nokey=1", str(file_path)],
                capture_output=True, text=True, timeout=30
            )
            return float(result.stdout.strip())
        except Exception:
            return 0.0

    def _extract_video_frames(self, file_path: Path, max_frames: int = 10) -> List[tuple]:
        """
        Extract key frames from video.
        Returns list of (frame_path, timestamp_seconds) tuples.
        """
        frames = []
        try:
            duration = self._get_video_duration(file_path)
            if duration <= 0:
                duration = 60  # Assume 60s if unknown

            interval = max(duration / max_frames, 2.0)  # At least 2s apart

            with tempfile.TemporaryDirectory() as tmpdir:
                for i in range(max_frames):
                    timestamp = i * interval
                    if timestamp >= duration:
                        break

                    frame_path = os.path.join(tmpdir, f"frame_{i:04d}.jpg")
                    result = subprocess.run(
                        ["ffmpeg", "-ss", str(timestamp), "-i", str(file_path),
                         "-vframes", "1", "-q:v", "2", frame_path, "-y"],
                        capture_output=True, timeout=30
                    )
                    if result.returncode == 0 and os.path.exists(frame_path):
                        # Copy to persistent temp location
                        persist_path = tempfile.NamedTemporaryFile(
                            suffix=".jpg", delete=False
                        ).name
                        import shutil
                        shutil.copy2(frame_path, persist_path)
                        frames.append((persist_path, timestamp))

        except Exception as e:
            print(f"[MediaExtractor] Frame extraction error: {e}")

        return frames

    def _extract_audio_transcript(self, file_path: Path) -> Optional[str]:
        """
        Extract audio and attempt transcription.
        Uses ffmpeg to extract audio, then sends to whisper if available.
        """
        try:
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
                audio_path = tmp.name

            result = subprocess.run(
                ["ffmpeg", "-i", str(file_path), "-vn", "-acodec", "pcm_s16le",
                 "-ar", "16000", "-ac", "1", audio_path, "-y"],
                capture_output=True, timeout=120
            )

            if result.returncode != 0 or not os.path.exists(audio_path):
                return None

            # Check if whisper is available
            try:
                import whisper
                model = whisper.load_model("base")
                result = model.transcribe(audio_path)
                os.unlink(audio_path)
                return result.get("text", "")
            except ImportError:
                pass

            # Try OpenAI Whisper API if provider available
            if self._vision_provider and hasattr(self._vision_provider, '_client'):
                try:
                    with open(audio_path, "rb") as audio_file:
                        transcript = self._vision_provider._client.audio.transcriptions.create(
                            model="whisper-1",
                            file=audio_file
                        )
                    os.unlink(audio_path)
                    return transcript.text
                except Exception:
                    pass

            os.unlink(audio_path)
            return None

        except Exception as e:
            print(f"[MediaExtractor] Audio extraction error: {e}")
            return None

    # ============================================================
    # Excel Extraction
    # ============================================================
    def _extract_excel(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract text content from Excel files."""
        if not OPENPYXL_AVAILABLE:
            return [{"content": f"[EXCEL: {file_path.name}] - openpyxl not installed.",
                      "metadata": {"source": str(file_path), "media_type": "excel", "error": "openpyxl_missing"}}]

        chunks = []
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows_data = []
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        rows_data.append(row_text)

                if rows_data:
                    sheet_content = "\n".join(rows_data)
                    # Chunk large sheets
                    if len(sheet_content) > 3000:
                        for i in range(0, len(sheet_content), 2500):
                            chunk_text = sheet_content[i:i+2500]
                            chunks.append({
                                "content": f"[EXCEL: {file_path.name} | Sheet: {sheet_name}]\n{chunk_text}",
                                "metadata": {
                                    "source": str(file_path),
                                    "media_type": "excel",
                                    "sheet_name": sheet_name,
                                    "file_name": file_path.name,
                                    "chunk_index": len(chunks)
                                }
                            })
                    else:
                        chunks.append({
                            "content": f"[EXCEL: {file_path.name} | Sheet: {sheet_name}]\n{sheet_content}",
                            "metadata": {
                                "source": str(file_path),
                                "media_type": "excel",
                                "sheet_name": sheet_name,
                                "file_name": file_path.name,
                                "chunk_index": len(chunks)
                            }
                        })

            wb.close()
        except Exception as e:
            chunks.append({
                "content": f"[EXCEL: {file_path.name}] - Error: {str(e)}",
                "metadata": {"source": str(file_path), "media_type": "excel", "error": str(e)}
            })

        return chunks if chunks else [{"content": f"[EXCEL: {file_path.name}] - Empty or unreadable.",
                                         "metadata": {"source": str(file_path), "media_type": "excel"}}]

    # ============================================================
    # CSV Extraction
    # ============================================================
    def _extract_csv(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract text content from CSV files."""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                rows = list(reader)

            if not rows:
                return [{"content": f"[CSV: {file_path.name}] - Empty file.",
                          "metadata": {"source": str(file_path), "media_type": "csv"}}]

            # Header + data
            header = rows[0] if rows else []
            data_rows = rows[1:] if len(rows) > 1 else []

            header_text = ", ".join(header)
            rows_text = "\n".join(
                ", ".join(str(cell) for cell in row) for row in data_rows
            )

            full_text = f"Headers: {header_text}\n{rows_text}"

            # Chunk large CSVs
            if len(full_text) > 3000:
                for i in range(0, len(full_text), 2500):
                    chunks.append({
                        "content": f"[CSV: {file_path.name}]\n{full_text[i:i+2500]}",
                        "metadata": {
                            "source": str(file_path),
                            "media_type": "csv",
                            "file_name": file_path.name,
                            "chunk_index": len(chunks)
                        }
                    })
            else:
                chunks.append({
                    "content": f"[CSV: {file_path.name}]\n{full_text}",
                    "metadata": {
                        "source": str(file_path),
                        "media_type": "csv",
                        "file_name": file_path.name,
                        "chunk_index": 0
                    }
                })

        except Exception as e:
            chunks.append({
                "content": f"[CSV: {file_path.name}] - Error: {str(e)}",
                "metadata": {"source": str(file_path), "media_type": "csv", "error": str(e)}
            })

        return chunks

    # ============================================================
    # Word Document Extraction
    # ============================================================
    def _extract_docx(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract text from Word documents."""
        if not PYTHON_DOCX_AVAILABLE:
            return [{"content": f"[DOCX: {file_path.name}] - python-docx not installed.",
                      "metadata": {"source": str(file_path), "media_type": "docx", "error": "python-docx_missing"}}]

        chunks = []
        try:
            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n\n".join(paragraphs)

            # Also extract tables
            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(" | ".join(cells))
                if table_rows:
                    full_text += "\n\n[TABLE]\n" + "\n".join(table_rows)

            # Chunk
            if len(full_text) > 3000:
                for i in range(0, len(full_text), 2500):
                    chunks.append({
                        "content": f"[DOCX: {file_path.name}]\n{full_text[i:i+2500]}",
                        "metadata": {
                            "source": str(file_path),
                            "media_type": "docx",
                            "file_name": file_path.name,
                            "chunk_index": len(chunks)
                        }
                    })
            else:
                chunks.append({
                    "content": f"[DOCX: {file_path.name}]\n{full_text}",
                    "metadata": {
                        "source": str(file_path),
                        "media_type": "docx",
                        "file_name": file_path.name,
                        "chunk_index": 0
                    }
                })

        except Exception as e:
            chunks.append({
                "content": f"[DOCX: {file_path.name}] - Error: {str(e)}",
                "metadata": {"source": str(file_path), "media_type": "docx", "error": str(e)}
            })

        return chunks if chunks else [{"content": f"[DOCX: {file_path.name}] - Empty document.",
                                         "metadata": {"source": str(file_path), "media_type": "docx"}}]

    # ============================================================
    # PowerPoint Extraction
    # ============================================================
    def _extract_pptx(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract text from PowerPoint presentations."""
        if not PYTHON_PPTX_AVAILABLE:
            return [{"content": f"[PPTX: {file_path.name}] - python-pptx not installed.",
                      "metadata": {"source": str(file_path), "media_type": "pptx", "error": "python-pptx_missing"}}]

        chunks = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    content = "\n".join(slide_text)
                    chunks.append({
                        "content": f"[PPTX: {file_path.name} | Slide {slide_num}]\n{content}",
                        "metadata": {
                            "source": str(file_path),
                            "media_type": "pptx",
                            "slide_number": slide_num,
                            "file_name": file_path.name,
                            "chunk_index": len(chunks)
                        }
                    })

        except Exception as e:
            chunks.append({
                "content": f"[PPTX: {file_path.name}] - Error: {str(e)}",
                "metadata": {"source": str(file_path), "media_type": "pptx", "error": str(e)}
            })

        return chunks if chunks else [{"content": f"[PPTX: {file_path.name}] - Empty presentation.",
                                         "metadata": {"source": str(file_path), "media_type": "pptx"}}]


# Singleton
_media_extractor: Optional[MediaExtractorService] = None


def get_media_extractor(vision_provider=None) -> MediaExtractorService:
    """Get the global media extractor instance."""
    global _media_extractor
    if _media_extractor is None:
        _media_extractor = MediaExtractorService(vision_provider=vision_provider)
    elif vision_provider and not _media_extractor._vision_provider:
        _media_extractor.set_vision_provider(vision_provider)
    return _media_extractor
